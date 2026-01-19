import os
import convertapi
from pptx import Presentation

# Initialize the conversion API with your key from the config
convertapi.api_credentials = os.getenv("CONVERTAPI_KEY")

def extract_text_slidewise(ppt_path):
    """
    Parses the PPTX file to extract text from every slide.
    This text is sent to Gemini so it knows what it is 'looking' at.
    """
    prs = Presentation(ppt_path)
    slides = []

    # Iterate through slides with a 1-based index for easy navigation
    for index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            # Only extract text from shapes that actually contain text frames
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        texts.append(p.text.strip())

        # Store the extracted text alongside the slide number
        slides.append({
            "slide_number": index,
            "text": " ".join(texts)
        })

    return slides

def convert_ppt_to_images(ppt_path, output_dir):
    """
    Converts the PPTX into high-quality JPG images.
    These images are shown in the 'Presentation View' on the frontend.
    """
    # Use ConvertAPI to transform the PPTX into a series of JPGs
    result = convertapi.convert("jpg", {"File": ppt_path})
    result.save_files(output_dir)

    # Gather all generated JPG files from the temporary directory
    images = [
        os.path.join(output_dir, f)
        for f in os.listdir(output_dir)
        if f.lower().endswith(".jpg")
    ]

    # CRITICAL: Sort by creation time to ensure slide 1 is first.
    # Without this, the slides might appear out of order on the frontend.
    images.sort(key=os.path.getctime)

    return images
